#!/usr/bin/env python3
"""Build the Rotana HR identity PowerPoint kit — native shapes, editable charts."""
from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "assets"
DIST = ROOT / "dist"
LOGO = ASSETS / "logo" / "rotana-logo-1400.png"
SPHERE = ASSETS / "logo" / "rotana-sphere.png"
WAVE = ASSETS / "logo" / "wave-motif.png"
ICON = ASSETS / "icons" / "png"
MASK = ASSETS / "masks"

FOREST = RGBColor(0x00, 0x51, 0x2F)
GROVE = RGBColor(0x00, 0x6C, 0x3C)
LEAF = RGBColor(0x2B, 0xA8, 0x6C)
INK = RGBColor(0x0C, 0x10, 0x0E)
INK2 = RGBColor(0x1A, 0x24, 0x20)
PAPER = RGBColor(0xF7, 0xF6, 0xF3)
MIST = RGBColor(0xE8, 0xE8, 0xE4)
STONE = RGBColor(0xC5, 0xC7, 0xC2)
SAGE = RGBColor(0xB8, 0xC0, 0xB0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BODY = RGBColor(0x3A, 0x40, 0x3C)
MUTED = RGBColor(0x6A, 0x70, 0x6C)

AR = "IBM Plex Sans Arabic"
AMIRI = "Amiri"
LATIN = "IBM Plex Sans"
MONO = "IBM Plex Mono"

W, H = Inches(13.333333), Inches(7.5)
MX, MY = Inches(0.62), Inches(0.42)


def rgb_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def line_only(shape, color: RGBColor, pt: float = 1.0) -> None:
    shape.fill.background()
    shape.line.color.rgb = color
    shape.line.width = Pt(pt)


def set_geom(shape, prst: str) -> None:
    spPr = shape._element.spPr
    geom = spPr.find(qn("a:prstGeom"))
    if geom is not None:
        geom.set("prst", prst)


def _ensure_typeface(rPr, tag: str, name: str) -> None:
    el = rPr.find(qn(f"a:{tag}"))
    if el is None:
        el = etree.SubElement(rPr, qn(f"a:{tag}"))
    el.set("typeface", name)


def style_run(run, name: str, size: float, color: RGBColor, bold: bool = False) -> None:
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    _ensure_typeface(rPr, "latin", name)
    _ensure_typeface(rPr, "ea", name)
    _ensure_typeface(rPr, "cs", name)


def p_rtl(paragraph, align: str = "r") -> None:
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("rtl", "1")
    pPr.set("algn", align)


def add_text(
    slide,
    l,
    t,
    w,
    h,
    text: str,
    *,
    size=18,
    color=INK,
    font=AR,
    bold=False,
    align="r",
    anchor=MSO_ANCHOR.TOP,
) -> None:
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    p = tf.paragraphs[0]
    p_rtl(p, align)
    run = p.add_run()
    run.text = text
    style_run(run, font, size, color, bold)
    return box


def add_lines(slide, l, t, w, h, lines: list[tuple[str, dict]], *, align="r") -> None:
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, kw) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p_rtl(p, kw.get("align", align))
        p.space_after = Pt(kw.get("after", 6))
        run = p.add_run()
        run.text = text
        style_run(
            run,
            kw.get("font", AR),
            kw.get("size", 16),
            kw.get("color", INK),
            kw.get("bold", False),
        )
    return box


def rect(slide, l, t, w, h, color: RGBColor):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    rgb_fill(sh, color)
    return sh


def oval(slide, l, t, w, h, color: RGBColor):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    rgb_fill(sh, color)
    return sh


def bg(slide, color: RGBColor) -> None:
    rect(slide, 0, 0, W, H, color)


def footer(slide, n: int, total: int = 24, dark: bool = False) -> None:
    c = SAGE if dark else MUTED
    add_text(slide, MX, Inches(7.12), Inches(5), Inches(0.28), "روتانا  ·  الموارد البشرية", size=11, color=c, font=AR)
    add_text(
        slide,
        Inches(10.4),
        Inches(7.12),
        Inches(2.3),
        Inches(0.28),
        f"{n:02d}  /  {total:02d}",
        size=11,
        color=c,
        font=MONO,
        align="l",
    )


def head(slide, idx: str, eyebrow: str, title: str) -> None:
    add_text(slide, MX, MY, Inches(0.7), Inches(0.3), idx, size=11, color=FOREST, font=MONO)
    add_text(slide, Inches(1.35), MY, Inches(11), Inches(0.28), eyebrow, size=11, color=FOREST, font=LATIN)
    add_text(slide, Inches(1.35), Inches(0.68), Inches(11.2), Inches(0.7), title, size=28, color=INK, font=AMIRI)


def icon(slide, name: str, l, t, s=Inches(0.38), white: bool = False) -> None:
    fn = ICON / (f"{name}-white.png" if white else f"{name}.png")
    if fn.exists():
        slide.shapes.add_picture(str(fn), l, t, s, s)


def picture(slide, path: Path, l, t, w, h, geom: str | None = None):
    pic = slide.shapes.add_picture(str(path), l, t, w, h)
    if geom:
        set_geom(pic, geom)
    return pic


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def s01(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, INK)
    rect(sl, 0, 0, W, Inches(0.05), FOREST)
    if LOGO.exists():
        sl.shapes.add_picture(str(LOGO), MX, Inches(0.48), Inches(3.4), Inches(0.78))
    add_text(sl, Inches(7.4), Inches(0.62), Inches(5.3), Inches(0.3), "HUMAN RESOURCES IDENTITY SYSTEM", size=11, color=SAGE, font=LATIN, align="l")
    add_text(sl, MX, Inches(2.15), Inches(11), Inches(0.35), "نظام العروض  ·  مرجع قابل للنسخ", size=14, color=LEAF, font=AR)
    add_text(sl, MX, Inches(2.55), Inches(11), Inches(2.1), "إدارة\nالموارد البشرية", size=60, color=WHITE, font=AMIRI)
    add_text(
        sl,
        MX,
        Inches(4.85),
        Inches(8.5),
        Inches(1.0),
        "قوالب، ماسكات، أرقام، أيقونات، وإدارات — تُنسخ كما هي وتُعدَّل كما تحتاجون.",
        size=16,
        color=SAGE,
        font=AR,
    )
    if WAVE.exists():
        sl.shapes.add_picture(str(WAVE), Inches(0.2), Inches(6.15), Inches(4.2), Inches(1.05))
    add_text(sl, MX, Inches(7.12), Inches(10), Inches(0.25), "روتانا   ·   هوية داخلية   ·   ١٦:٩", size=12, color=STONE, font=AR)
    return sl


def s02(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, PAPER)
    head(sl, "٠٢", "طريقة الاستخدام", "انسخ الشريحة. استبدل الصورة. غيّر الرقم.")
    items = [
        ("١", "النسخ من هذا الملف", "كل شريحة تخطيط جاهز. انسخ الشريحة إلى عرضك، أو انسخ العنصر وحده."),
        ("٢", "الاستبدال لا إعادة الرسم", "الصور داخل ماسكات. Change Picture. الأيقونات ملفات مستقلة. الألوان لها أكواد."),
        ("٣", "لا تختَرع لونًا جديدًا", "الأخضر من الشعار، والرمادي الفاتح هو الورق. الفخامة في الفراغ والخط."),
    ]
    for i, (n, t, b) in enumerate(items):
        x = MX + Inches(i * 4.1)
        add_text(sl, x, Inches(2.0), Inches(3.7), Inches(0.7), n, size=40, color=FOREST, font=AMIRI)
        add_text(sl, x, Inches(2.75), Inches(3.7), Inches(0.4), t, size=18, color=INK, font=AR, bold=True)
        add_text(sl, x, Inches(3.25), Inches(3.7), Inches(1.6), b, size=14, color=BODY, font=AR)
    add_text(sl, MX, Inches(6.55), Inches(12), Inches(0.35), "الشريحة ٠١ و٢٤ للغلاف. الباقي مادة عمل يومي.", size=13, color=MUTED, font=AR)
    footer(sl, 2)
    return sl


def s03(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, PAPER)
    head(sl, "٠٣", "نظام اللون", "من الشعار، ومعه الرمادي الفاتح.")
    swatches = [
        (FOREST, WHITE, "Forest", "#00512F", "RGB 0 · 81 · 47", "الأساسي — عناوين وقواعد"),
        (GROVE, WHITE, "Grove", "#006C3C", "RGB 0 · 108 · 60", "وسط الكرة — رسوم"),
        (LEAF, INK, "Leaf", "#2BA86C", "RGB 43 · 168 · 108", "لمسة فقط"),
        (MIST, INK, "Mist", "#E8E8E4", "RGB 232 · 232 · 228", "الرمادي الفاتح — حقول"),
        (INK, WHITE, "Ink", "#0C100E", "RGB 12 · 16 · 14", "أغلفة وأقسام"),
        (STONE, INK, "Stone", "#C5C7C2", "RGB 197 · 199 · 194", "خطوط رفيعة"),
    ]
    for i, (bgc, tc, name, hex_, rgb, use) in enumerate(swatches):
        x = MX + Inches(i * 2.05)
        rect(sl, x, Inches(1.7), Inches(1.92), Inches(4.85), bgc)
        add_text(sl, x + Inches(0.12), Inches(4.05), Inches(1.7), Inches(0.3), name, size=14, color=tc, font=LATIN)
        add_text(sl, x + Inches(0.12), Inches(4.38), Inches(1.7), Inches(0.28), hex_, size=12, color=tc, font=MONO)
        add_text(sl, x + Inches(0.12), Inches(4.66), Inches(1.7), Inches(0.25), rgb, size=9, color=tc, font=MONO)
        add_text(sl, x + Inches(0.12), Inches(5.05), Inches(1.7), Inches(1.1), use, size=12, color=tc, font=AR)
    footer(sl, 3)
    return sl


def s04(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, PAPER)
    head(sl, "٠٤", "تطبيق اللون", "أخضر للقرار. رمادي للهدوء. حبر للوقفة.")
    rect(sl, MX, Inches(1.75), Inches(4.0), Inches(2.55), INK)
    add_text(sl, MX + Inches(0.25), Inches(3.45), Inches(3.5), Inches(0.55), "غلاف / فاصل\nInk + Forest", size=16, color=WHITE, font=AR)
    rect(sl, Inches(4.82), Inches(1.75), Inches(4.0), Inches(2.55), MIST)
    add_text(sl, Inches(5.07), Inches(3.45), Inches(3.5), Inches(0.55), "محتوى يومي\nPaper + Mist", size=16, color=INK, font=AR)
    rect(sl, Inches(8.72), Inches(1.75), Inches(4.0), Inches(2.55), FOREST)
    add_text(sl, Inches(8.97), Inches(3.45), Inches(3.5), Inches(0.55), "رقم واحد مهم\nForest field", size=16, color=WHITE, font=AR)
    rules = [
        "لا تضع تدرجًا على النص.",
        "لا تستخدم الذهبي أو النيون — ليسا من الشعار.",
        "Leaf للبيانات الموجبة فقط، بمساحة صغيرة.",
        "الحدود Stone بسمك ١px، لا ظلال تحت البطاقات.",
    ]
    for i, r in enumerate(rules):
        col, row = i % 2, i // 2
        x = MX + Inches(col * 6.2)
        y = Inches(4.6) + Inches(row * 0.55)
        oval(sl, x, y + Inches(0.1), Inches(0.12), Inches(0.12), FOREST)
        add_text(sl, x + Inches(0.28), y, Inches(5.6), Inches(0.4), r, size=15, color=BODY, font=AR)
    footer(sl, 4)
    return sl


def s05(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, PAPER)
    head(sl, "٠٥", "نظام الخط", "عربي أولاً. لاتيني يمشي معه.")
    add_text(sl, MX, Inches(1.7), Inches(4.1), Inches(0.28), "DISPLAY — AMIRI", size=11, color=FOREST, font=LATIN)
    add_text(sl, MX, Inches(2.05), Inches(4.2), Inches(2.2), "المواهب تُدار\nكما تُدار العلامة.", size=28, color=INK, font=AMIRI)
    add_text(sl, MX, Inches(4.4), Inches(4.2), Inches(0.4), "عناوين الأغلفة والفواصل  ·  ٦٤–٩٢", size=12, color=MUTED, font=AR)

    add_text(sl, Inches(5.1), Inches(1.7), Inches(4.2), Inches(0.28), "UI — IBM PLEX SANS ARABIC", size=11, color=FOREST, font=LATIN)
    add_text(sl, Inches(5.1), Inches(2.1), Inches(4.2), Inches(1.0), "الاستقطاب، التهيئة، الأداء، المزايا، الثقافة.", size=18, color=INK, font=AR, bold=True)
    add_text(sl, Inches(5.1), Inches(3.2), Inches(4.2), Inches(1.2), "الجسم ١٨–٢٠. الشروحات ١٥. التسميات ١٢. لا تضغط السطر.", size=14, color=BODY, font=AR)

    add_text(sl, Inches(9.5), Inches(1.7), Inches(3.3), Inches(0.28), "LATIN — IBM PLEX SANS", size=11, color=FOREST, font=LATIN)
    add_text(sl, Inches(9.5), Inches(2.15), Inches(3.3), Inches(0.7), "Rotana HR\nHeadcount · eNPS", size=16, color=INK, font=LATIN)
    add_text(sl, Inches(9.5), Inches(3.2), Inches(3.3), Inches(0.28), "CODES — IBM PLEX MONO", size=11, color=FOREST, font=LATIN)
    add_text(sl, Inches(9.5), Inches(3.55), Inches(3.3), Inches(0.5), "#00512F", size=18, color=FOREST, font=MONO)
    footer(sl, 5)
    return sl


def s06(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, PAPER)
    head(sl, "٠٦", "الشعار", "لا يُعاد رسمه. يُوضع كما هو.")
    rect(sl, MX, Inches(1.7), Inches(5.5), Inches(4.15), INK)
    if LOGO.exists():
        sl.shapes.add_picture(str(LOGO), Inches(1.15), Inches(3.05), Inches(4.4), Inches(1.0))
    add_text(sl, MX, Inches(5.55), Inches(5.5), Inches(0.28), "على الحبر", size=12, color=SAGE, font=AR, align="ctr")

    rect(sl, Inches(6.35), Inches(1.7), Inches(3.6), Inches(4.15), MIST)
    if LOGO.exists():
        sl.shapes.add_picture(str(LOGO), Inches(6.65), Inches(3.05), Inches(3.0), Inches(0.7))
    add_text(sl, Inches(6.35), Inches(5.55), Inches(3.6), Inches(0.28), "على الرمادي الفاتح", size=12, color=BODY, font=AR, align="ctr")

    rect(sl, Inches(10.15), Inches(1.7), Inches(2.55), Inches(4.15), WHITE)
    line_only(
        sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.15), Inches(1.7), Inches(2.55), Inches(4.15)),
        MIST,
        1,
    )
    if SPHERE.exists():
        sl.shapes.add_picture(str(SPHERE), Inches(10.4), Inches(2.55), Inches(2.05), Inches(2.15))
    add_text(sl, Inches(10.15), Inches(5.55), Inches(2.55), Inches(0.28), "الكرة — أيقونة", size=12, color=BODY, font=AR, align="ctr")
    add_text(sl, MX, Inches(6.5), Inches(12), Inches(0.35), "مساحة الأمان = ارتفاع حرف r. لا ظل، ولا إعادة تلوين للوردمارك.", size=13, color=MUTED, font=AR)
    footer(sl, 6)
    return sl


def s07(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, PAPER)
    head(sl, "٠٧", "ماسكات الصورة", "شكل ثابت. صورة تُستبدل.")
    files = [
        (MASK / "circle-01.png", "دائري · قيادات", True),
        (MASK / "arch-01.png", "قوس · بورتريه", False),
        (MASK / "rounded-01.png", "مستطيل ناعم", False),
        (MASK / "wave-01.png", "موجة روتانا", False),
    ]
    for i, (path, cap, circ) in enumerate(files):
        x = MX + Inches(i * 3.15)
        if path.exists():
            pic = sl.shapes.add_picture(str(path), x, Inches(1.65), Inches(2.85), Inches(4.35))
            if circ:
                set_geom(pic, "ellipse")
        add_text(sl, x, Inches(6.1), Inches(2.85), Inches(0.4), cap, size=13, color=BODY, font=AR, align="ctr")
    add_text(sl, MX, Inches(6.55), Inches(12), Inches(0.3), "Picture Format → Change Picture. الماسك يبقى.", size=13, color=MUTED, font=AR)
    footer(sl, 7)
    return sl


ICON_LABELS = [
    ("people", "فريق"),
    ("hire", "استقطاب"),
    ("onboard", "تهيئة"),
    ("train", "تدريب"),
    ("performance", "أداء"),
    ("benefits", "مزايا"),
    ("relations", "علاقات"),
    ("culture", "ثقافة"),
    ("wellness", "رفاه"),
    ("payroll", "رواتب"),
    ("policy", "سياسات"),
    ("talent", "مواهب"),
    ("org", "هيكل"),
    ("calendar", "تقويم"),
    ("chart", "مؤشر"),
    ("target", "هدف"),
    ("handshake", "شراكة"),
    ("shield", "التزام"),
    ("globe", "عالمي"),
    ("music", "موسيقى"),
    ("film", "سينما"),
    ("tv", "تلفزيون"),
    ("megaphone", "إعلام"),
    ("building", "مقر"),
    ("document", "مستند"),
    ("clock", "زمن"),
    ("pin", "موقع"),
    ("message", "رسالة"),
]


def s08(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, PAPER)
    head(sl, "٠٨", "مكتبة الأيقونات", "خط واحد، نهايات مدوّرة — بلغة الشعار.")
    for i, (name, label) in enumerate(ICON_LABELS):
        col, row = i % 7, i // 7
        x = MX + Inches(col * 1.78)
        y = Inches(1.7) + Inches(row * 1.22)
        icon(sl, name, x + Inches(0.55), y, Inches(0.42))
        add_text(sl, x, y + Inches(0.5), Inches(1.55), Inches(0.32), label, size=12, color=BODY, font=AR, align="ctr")
    footer(sl, 8)
    return sl


def s09(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, PAPER)
    head(sl, "٠٩", "إنفوجرافيك الأرقام", "رقم كبير. شرح صغير. بلا صندوق.")
    kpis = [
        ("1,240", "إجمالي القوة", "Headcount · Q2"),
        ("87%", "نسبة التهيئة", "Onboarding complete"),
        ("24", "يومًا للتوظيف", "Time to hire"),
        ("+6.2", "eNPS", "مقابل الربع السابق"),
    ]
    for i, (n, ar, en) in enumerate(kpis):
        x = MX + Inches(i * 3.15)
        add_text(sl, x, Inches(1.7), Inches(2.95), Inches(1.15), n, size=44, color=FOREST, font=AMIRI)
        add_text(sl, x, Inches(2.85), Inches(2.95), Inches(0.35), ar, size=16, color=INK, font=AR, bold=True)
        add_text(sl, x, Inches(3.2), Inches(2.95), Inches(0.3), en, size=12, color=MUTED, font=LATIN)
    rect(sl, MX, Inches(3.7), Inches(12.1), Emu(12700), STONE)
    bars = [("استبقاء السنة الأولى", 0.78, "78%"), ("اكتمال التقييمات", 0.64, "64%"), ("شواغر حرجة مغلقة", 0.41, "41%")]
    for i, (lab, frac, val) in enumerate(bars):
        y = Inches(4.05) + Inches(i * 0.72)
        add_text(sl, MX, y, Inches(3.3), Inches(0.3), lab, size=14, color=BODY, font=AR)
        rect(sl, Inches(4.1), y + Inches(0.08), Inches(7.4), Inches(0.16), MIST)
        rect(sl, Inches(4.1) + Inches(7.4 * (1 - frac)), y + Inches(0.08), Inches(7.4 * frac), Inches(0.16), FOREST)
        add_text(sl, Inches(11.6), y, Inches(1.1), Inches(0.3), val, size=13, color=FOREST, font=MONO, align="l")
    footer(sl, 9)
    return sl


def _recolor_chart(chart, series_colors: list[RGBColor]) -> None:
    for i, color in enumerate(series_colors):
        try:
            ser = chart.series[i]
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = color
        except Exception:
            pass


def s10(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, PAPER)
    head(sl, "١٠", "تكوين القوة", "دائرة القرار، وأعمدة المقارنة — عدّل الأرقام من الشارت.")
    donut = CategoryChartData()
    donut.categories = ["تشغيل", "دعم", "قيادة"]
    donut.add_series("تكوين", (62, 23, 15))
    chart = sl.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, MX, Inches(1.65), Inches(4.6), Inches(4.6), donut
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    try:
        plot = chart.plots[0]
        plot.has_data_labels = True
        series = chart.series[0]
        for pt, col in zip(series.points, [FOREST, GROVE, MIST]):
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = col
    except Exception:
        pass

    bars = CategoryChartData()
    bars.categories = ["موسيقى", "تلفزيون", "سينما", "رقمي", "شركات"]
    bars.add_series("حصة", (82, 64, 48, 36, 28))
    bchart = sl.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(6.4), Inches(1.85), Inches(6.4), Inches(4.4), bars
    ).chart
    bchart.has_legend = False
    try:
        bchart.series[0].format.fill.solid()
        bchart.series[0].format.fill.fore_color.rgb = FOREST
    except Exception:
        pass
    footer(sl, 10)
    return sl


def s11(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, PAPER)
    head(sl, "١١", "إدارات روتانا", "وحدات العمل — انسخ البطاقة إلى عرض القسم.")
    depts = [
        "روتانا موسيقى",
        "روتانا خليجية",
        "روتانا سينما",
        "روتانا دراما",
        "روتانا كوميدي",
        "روتانا زمان",
        "روتانا كلاسيك",
        "روتانا كيدز",
        "المنصات الرقمية",
        "التسويق والعلامات",
        "المبيعات والشراكات",
        "الإنتاج",
        "التوزيع والحقوق",
        "المالية",
        "الشؤون القانونية",
        "تقنية المعلومات",
        "الاتصال المؤسسي",
        "الموارد البشرية",
    ]
    for i, name in enumerate(depts):
        col, row = i % 3, i // 3
        x = MX + Inches(col * 4.15)
        y = Inches(1.65) + Inches(row * 0.78)
        box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.95), Inches(0.62))
        rgb_fill(box, WHITE)
        box.line.color.rgb = STONE
        box.line.width = Pt(1)
        add_text(sl, x + Inches(0.18), y + Inches(0.14), Inches(3.6), Inches(0.36), name, size=15, color=INK, font=AR)
    footer(sl, 11)
    return sl


HR_UNITS = [
    ("hire", "الاستقطاب", "Talent Acquisition"),
    ("onboard", "التهيئة", "Onboarding"),
    ("train", "التدريب والتطوير", "L&D"),
    ("performance", "إدارة الأداء", "Performance"),
    ("benefits", "التعويضات والمزايا", "C&B"),
    ("relations", "علاقات الموظفين", "ER"),
    ("culture", "الثقافة والرفاه", "Culture"),
    ("org", "تخطيط القوة", "Workforce"),
    ("policy", "السياسات والالتزام", "Compliance"),
    ("payroll", "الرواتب", "Payroll"),
    ("talent", "المواهب والخلافة", "Talent"),
    ("wellness", "الصحة المهنية", "Wellbeing"),
]


def s12(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, PAPER)
    head(sl, "١٢", "وحدات الموارد البشرية", "كل وظيفة لها أيقونة — انسخ الصف.")
    for i, (ic, ar, en) in enumerate(HR_UNITS):
        col, row = i % 3, i // 3
        x = MX + Inches(col * 4.15)
        y = Inches(1.65) + Inches(row * 1.18)
        icon(sl, ic, x, y + Inches(0.18), Inches(0.4))
        add_text(sl, x + Inches(0.55), y + Inches(0.12), Inches(3.3), Inches(0.32), ar, size=16, color=INK, font=AR, bold=True)
        add_text(sl, x + Inches(0.55), y + Inches(0.44), Inches(3.3), Inches(0.28), en, size=11, color=MUTED, font=LATIN)
        rect(sl, x, y + Inches(0.9), Inches(3.85), Emu(6350), MIST)
    footer(sl, 12)
    return sl


def s13(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, INK)
    add_text(sl, MX, Inches(2.15), Inches(4), Inches(0.3), "٠١", size=14, color=LEAF, font=MONO)
    add_text(sl, MX, Inches(2.55), Inches(11), Inches(2.0), "الاستقطاب\nوالاختيار", size=54, color=WHITE, font=AMIRI)
    add_text(sl, MX, Inches(4.75), Inches(8), Inches(0.35), "TALENT ACQUISITION", size=13, color=SAGE, font=LATIN)
    if WAVE.exists():
        sl.shapes.add_picture(str(WAVE), Inches(0.2), Inches(6.15), Inches(4.2), Inches(1.05))
    footer(sl, 13, dark=True)
    return sl


def s14(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, PAPER)
    head(sl, "١٤", "جدول الأعمال", "اجتماع قيادة الموارد البشرية")
    items = [
        ("٠١", "نبض القوة العاملة", "١٢ د"),
        ("٠٢", "الشواغر الحرجة والمواهب", "١٨ د"),
        ("٠٣", "مسار التهيئة للربع", "١٠ د"),
        ("٠٤", "الثقافة ومؤشر eNPS", "١٢ د"),
        ("٠٥", "قرارات ومالكو المتابعة", "٨ د"),
    ]
    for i, (n, t, d) in enumerate(items):
        y = Inches(1.7) + Inches(i * 0.9)
        add_text(sl, MX, y + Inches(0.22), Inches(0.7), Inches(0.35), n, size=13, color=FOREST, font=MONO)
        add_text(sl, Inches(1.5), y + Inches(0.16), Inches(9.2), Inches(0.45), t, size=22, color=INK, font=AR)
        add_text(sl, Inches(11.3), y + Inches(0.22), Inches(1.4), Inches(0.35), d, size=14, color=MUTED, font=AR, align="l")
        rect(sl, MX, y + Inches(0.78), Inches(12.1), Emu(12700), STONE)
    footer(sl, 14)
    return sl


def s15(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, PAPER)
    head(sl, "١٥", "تخطيط عنوان + نص", "مبدأ واحد في الشريحة.")
    add_text(
        sl,
        MX,
        Inches(2.1),
        Inches(12.0),
        Inches(2.8),
        "نُوظّف للعلامة قبل أن نُوظّف للوظيفة. المرشح الذي لا يستطيع تمثيل روتانا في غرفة لاحقة، لا يُعوَّض بسرعة إجراءات التوظيف. هذا التخطيط لنص سياسات، أو لمقدمة عرض، أو لقرار يُقرأ بصوت عالٍ.",
        size=22,
        color=INK,
        font=AR,
    )
    add_text(sl, MX, Inches(5.2), Inches(12), Inches(0.5), "انسخ النص. غيّر الجملة. أبقِ الحجم والهامش.", size=15, color=MUTED, font=AR)
    footer(sl, 15)
    return sl


def s16(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, PAPER)
    head(sl, "١٦", "عمودان", "ما نُبقي عليه / ما نُغيّره.")
    add_text(sl, MX, Inches(1.7), Inches(5.6), Inches(0.5), "نُبقي", size=28, color=FOREST, font=AMIRI)
    keep = [
        "مقابلة الكفاءة الثقافية مع كل عرض.",
        "مسار تهيئة ٣٠–٦٠–٩٠ يومًا.",
        "مالك واحد لكل شاغر حرج.",
    ]
    change = [
        "زمن الاعتماد الداخلي فوق خمس طبقات.",
        "الوصف الوظيفي كمستند جامد.",
        "قياس التدريب بالحضور لا بالأثر.",
    ]
    for i, t in enumerate(keep):
        y = Inches(2.4) + Inches(i * 1.05)
        rect(sl, MX, y + Inches(0.85), Inches(5.6), Emu(12700), MIST)
        add_text(sl, MX, y, Inches(5.6), Inches(0.8), t, size=16, color=INK, font=AR)
    add_text(sl, Inches(7.3), Inches(1.7), Inches(5.6), Inches(0.5), "نُغيّر", size=28, color=FOREST, font=AMIRI)
    for i, t in enumerate(change):
        y = Inches(2.4) + Inches(i * 1.05)
        rect(sl, Inches(7.3), y + Inches(0.85), Inches(5.6), Emu(12700), MIST)
        add_text(sl, Inches(7.3), y, Inches(5.6), Inches(0.8), t, size=16, color=INK, font=AR)
    footer(sl, 16)
    return sl


def s17(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, PAPER)
    head(sl, "١٧", "ثلاث ركائز", "العام التشغيلي للموارد البشرية.")
    pillars = [
        ("people", "ناس", "استقطاب يطابق صوت روتانا، وتهيئة لا تترك أحدًا في الأسبوع الأول وحيدًا."),
        ("performance", "أداء", "أهداف واضحة، مراجعة نصف سنوية، وخلافة للمقاعد التي لا تحتمل فراغًا."),
        ("culture", "ثقافة", "مكان يُنتج فيه العمل الإبداعي دون فوضى إجراءات، وبلا تساهل في الاحترام."),
    ]
    for i, (ic, t, b) in enumerate(pillars):
        x = MX + Inches(i * 4.15)
        icon(sl, ic, x, Inches(1.85), Inches(0.48))
        add_text(sl, x, Inches(2.5), Inches(3.8), Inches(0.55), t, size=28, color=INK, font=AMIRI)
        add_text(sl, x, Inches(3.2), Inches(3.8), Inches(2.0), b, size=15, color=BODY, font=AR)
    footer(sl, 17)
    return sl


def s18(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, PAPER)
    head(sl, "١٨", "مسار من أربع محطات", "من الشاغر إلى اليوم التسعين.")
    rect(sl, Inches(1.6), Inches(2.55), Inches(10.1), Emu(12700), STONE)
    steps = [
        ("١", "تعريف", "الوصف، النطاق، ومالك القرار."),
        ("٢", "اختيار", "غربلة، مقابلة، ومحاكاة عمل."),
        ("٣", "عرض", "حزمة، تاريخ، وتوقيع واحد."),
        ("٤", "اندماج", "٣٠–٦٠–٩٠ بقياس لا بانطباع."),
    ]
    for i, (n, t, b) in enumerate(steps):
        x = MX + Inches(i * 3.15)
        oval(sl, x + Inches(1.05), Inches(2.32), Inches(0.52), Inches(0.52), FOREST)
        add_text(sl, x + Inches(1.05), Inches(2.38), Inches(0.52), Inches(0.42), n, size=16, color=WHITE, font=AMIRI, align="ctr")
        add_text(sl, x, Inches(3.15), Inches(2.9), Inches(0.45), t, size=20, color=INK, font=AR, bold=True, align="ctr")
        add_text(sl, x, Inches(3.65), Inches(2.9), Inches(1.1), b, size=14, color=BODY, font=AR, align="ctr")
    footer(sl, 18)
    return sl


def s19(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, PAPER)
    head(sl, "١٩", "خط زمني", "سنة الثقافة والأداء.")
    items = [
        ("Q1", "إطلاق الأهداف وتهيئة الدفعة الشتوية"),
        ("Q2", "مراجعة منتصف العام ومسار الخلافة"),
        ("Q3", "مسح eNPS وبرامج القادة"),
        ("Q4", "معايرة الأداء وتخطيط القوة للعام التالي"),
    ]
    for i, (q, t) in enumerate(items):
        y = Inches(1.75) + Inches(i * 1.15)
        add_text(sl, MX, y + Inches(0.2), Inches(1.3), Inches(0.4), q, size=18, color=FOREST, font=LATIN, bold=True)
        add_text(sl, Inches(2.2), y + Inches(0.16), Inches(10.4), Inches(0.5), t, size=22, color=INK, font=AR)
        rect(sl, MX, y + Inches(0.85), Inches(12.1), Emu(12700), STONE)
    footer(sl, 19)
    return sl


def s20(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, PAPER)
    head(sl, "٢٠", "فريق القيادة", "استبدل الصورة. أبقِ الإطار.")
    roles = [
        ("circle-01.png", "الاسم", "رئيس الموارد البشرية"),
        ("circle-02.png", "الاسم", "الاستقطاب"),
        ("circle-03.png", "الاسم", "التطوير والأداء"),
        ("circle-04.png", "الاسم", "المزايا والثقافة"),
    ]
    for i, (fn, name, role) in enumerate(roles):
        x = MX + Inches(i * 3.15)
        path = MASK / fn
        if path.exists():
            pic = sl.shapes.add_picture(str(path), x + Inches(0.25), Inches(1.7), Inches(2.45), Inches(2.45))
            set_geom(pic, "ellipse")
        add_text(sl, x, Inches(4.35), Inches(2.95), Inches(0.35), name, size=16, color=INK, font=AR, bold=True, align="ctr")
        add_text(sl, x, Inches(4.7), Inches(2.95), Inches(0.35), role, size=13, color=MUTED, font=AR, align="ctr")
    footer(sl, 20)
    return sl


def s21(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, INK)
    add_text(sl, MX, Inches(1.4), Inches(3), Inches(1.1), "«", size=80, color=FOREST, font=AMIRI)
    add_text(
        sl,
        MX,
        Inches(2.5),
        Inches(11.8),
        Inches(2.4),
        "العلامة تُسمَع في الأغنية، وتُرى فيمن يفتح الباب في اليوم الأول.",
        size=32,
        color=WHITE,
        font=AMIRI,
    )
    add_text(sl, MX, Inches(5.2), Inches(10), Inches(0.35), "ROTANA  ·  HUMAN RESOURCES", size=12, color=SAGE, font=LATIN)
    footer(sl, 21, dark=True)
    return sl


def s22(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, PAPER)
    head(sl, "٢٢", "مقارنة مسارات", "قبل / بعد إعادة تصميم التهيئة.")
    rows = [
        ("", "سابقًا", "الآن"),
        ("أول يوم", "أوراق وتوقيع", "لقاء فريق + مشغّل العمل"),
        ("الأسبوع ١", "بريد تعريفي", "شريك تهيئة مسمّى"),
        ("اليوم ٣٠", "لا قياس", "مراجعة ٣٠ يومًا"),
        ("اليوم ٩٠", "انطباع المدير", "ثلاثة مؤشرات أثر"),
    ]
    table = sl.shapes.add_table(len(rows), 3, MX, Inches(1.75), Inches(12.1), Inches(4.4)).table
    table.columns[0].width = Inches(2.4)
    table.columns[1].width = Inches(4.85)
    table.columns[2].width = Inches(4.85)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p_rtl(p, "r")
            run = p.add_run()
            run.text = val
            color = FOREST if r == 0 else INK
            font = LATIN if r == 0 else AR
            size = 12 if r == 0 else 15
            style_run(run, font, size, color, bold=(r == 0 or c == 0))
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = PAPER if r % 2 == 0 else RGBColor(0xF0, 0xEF, 0xEC)
    footer(sl, 22)
    return sl


def s23(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, PAPER)
    add_text(sl, MX, Inches(1.8), Inches(5.6), Inches(0.3), "تخطيط صورة + نص", size=12, color=FOREST, font=LATIN)
    add_text(sl, MX, Inches(2.2), Inches(5.6), Inches(1.3), "المكان يسبق الشريحة.", size=32, color=INK, font=AMIRI)
    add_text(
        sl,
        MX,
        Inches(3.6),
        Inches(5.6),
        Inches(1.8),
        "هذا الماسك للموجة. ضع صورة استوديو أو بروفة أو فريق. النص على الورق، والصورة تأخذ القرار البصري.",
        size=16,
        color=BODY,
        font=AR,
    )
    add_text(sl, MX, Inches(5.6), Inches(5.6), Inches(0.4), "Change Picture على الموجة فقط.", size=13, color=MUTED, font=LATIN)
    wave = MASK / "wave-01.png"
    if wave.exists():
        sl.shapes.add_picture(str(wave), Inches(7.15), Inches(0.55), Inches(5.7), Inches(6.4))
    footer(sl, 23)
    return sl


def s24(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, INK)
    if LOGO.exists():
        sl.shapes.add_picture(str(LOGO), MX, Inches(2.15), Inches(4.0), Inches(0.92))
    add_text(sl, MX, Inches(3.35), Inches(10), Inches(1.1), "شكرًا", size=64, color=WHITE, font=AMIRI)
    add_text(sl, MX, Inches(4.55), Inches(10), Inches(0.4), "إدارة الموارد البشرية", size=18, color=SAGE, font=AR)
    add_text(sl, MX, Inches(6.4), Inches(10), Inches(0.3), "انسخ  ·  استبدل  ·  لا تُعدّل الأصل", size=13, color=STONE, font=AR)
    footer(sl, 24, dark=True)
    return sl


def build() -> Path:
    DIST.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    for fn in (
        s01, s02, s03, s04, s05, s06, s07, s08, s09, s10,
        s11, s12, s13, s14, s15, s16, s17, s18, s19, s20,
        s21, s22, s23, s24,
    ):
        fn(prs)
    out = DIST / "Rotana_HR_Identity_Kit.pptx"
    prs.save(out)
    print("wrote", out, "slides", len(prs.slides))
    return out


if __name__ == "__main__":
    build()
